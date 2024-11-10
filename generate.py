import logging
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import matplotlib.pyplot as plt
import plotly.express as px
from langchain.llms import HuggingFaceEndpoint
import torch
import os
from PIL import Image

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set up Hugging Face API token
HUGGINGFACE_API_TOKEN = ""  # Replace with your Hugging Face API Token

# Force PyTorch to use CPU
device = torch.device("cpu")

# Set up Streamlit page configuration
st.set_page_config(page_title="LLM-Powered Presentation Generator", page_icon="📊")
st.header("Automatic LLM-Powered Presentation Generator 📊")

# Sidebar for user input
with st.sidebar:
    st.title("Create Your Presentation")
    topic = st.text_input("Enter the Topic of your Presentation:")

    # Option to generate charts for the presentation
    generate_charts = st.checkbox("Generate Pie Chart for Content")
    bar_chart = st.checkbox("Generate Bar Chart for Content")

    # Load templates and display them as options
    templates_folder = "templates/"
    templates = [f for f in os.listdir(templates_folder) if f.endswith('.pptx')]
    template_names = [t.replace('.pptx', '') for t in templates]
    
    # Display template preview images
    selected_template_name = st.radio("Select a Template", template_names)

    # Show the selected template's preview image
    preview_image_path = os.path.join(templates_folder, f"{selected_template_name}_preview.png")
    try:
        img = Image.open(preview_image_path)
        st.image(img, caption=f"Preview of {selected_template_name}", use_column_width=True)
    except Exception as e:
        st.warning("No preview available for the selected template.")

# Load the LLM for Presentation Generation
llm = HuggingFaceEndpoint(
    repo_id="mistralai/Mistral-Nemo-Instruct-2407",
    max_new_tokens=512,
    top_k=10,
    top_p=0.95,
    typical_p=0.95,
    temperature=0.01,
    repetition_penalty=1.03,
    huggingfacehub_api_token=HUGGINGFACE_API_TOKEN
)

# Function to generate presentation content using the Hugging Face model
def generate_presentation_content(topic, progress_bar):
    """Generates presentation content using the Hugging Face LLM."""
    try:
        # Update progress bar
        progress_bar.progress(30)  # Mark 30% as done (content generation starts)
        
        # Prepare the prompt for the LLM
        prompt = f"""
        Create a PowerPoint presentation on the topic of: {topic}. 
        The presentation should include the following sections:
        Introduction: Brief overview or background of the topic.
        Main Points: Key ideas or arguments related to the topic.
        Supporting Evidence or Examples: Data, case studies, or relevant details that support the main points.
        Conclusion: A summary of the key takeaways or a call to action.
        Feel free to adjust the content to make it relevant to the topic. 
        Add visuals, charts, or graphs where appropriate to enhance understanding."""
        
        # Invoke the LLM to generate content
        response = llm.invoke(prompt)
        
        # Mark 100% completion after generating content
        progress_bar.progress(100)
        
        return response
    except Exception as e:
        logger.error(f"Error while generating content: {e}")
        return "Error generating content."

# Function to safely get a placeholder from a slide at a given index
def get_placeholder(slide, index):
    """Safely get the placeholder at a specific index in the slide."""
    try:
        # Attempt to access the placeholder using index-based access (which might raise a KeyError)
        return slide.placeholders[index]
    except KeyError as e:
        # Catch KeyError, which happens if the placeholder index does not exist
        logger.warning(f"KeyError: Placeholder index {index} not found in this slide. Error: {e}")
        return None
    except IndexError as e:
        # Catch IndexError, which happens if the index is out of bounds
        logger.warning(f"IndexError: Placeholder index {index} not found in this slide. Error: {e}")
        return None
    except Exception as e:
        # Catch any other unforeseen errors
        logger.error(f"Unexpected error accessing placeholder index {index}: {e}")
        return None

# Function to create PowerPoint presentation from selected template
def create_presentation_from_template(content, topic, template_file, progress_bar):
    """Creates a PowerPoint presentation from the selected template."""
    prs = Presentation(template_file)
    
    # Title Slide (if necessary, you can customize it per template)
    slide_layout = prs.slide_layouts[0]  # 0 for title slide (most templates have this as the first slide layout)
    slide = prs.slides[0]  # Use the first slide in the template
    title = slide.shapes.title
    
    # Safely access subtitle placeholder (index 1)
    subtitle = get_placeholder(slide, 1)  # Ensure placeholder exists
    title.text = topic
    
    if subtitle:  # Check if the subtitle placeholder exists
        subtitle.text = "Generated by LLM"
    else:
        logger.warning(f"No subtitle placeholder found in the first slide of the template {template_file}")

    # Add a content slide with generated content
    slide_layout = prs.slide_layouts[1]  # 1 for title and content
    slide = prs.slides.add_slide(slide_layout)

    # Safely access the title and content placeholders (index 1)
    title = slide.shapes.title
    body = get_placeholder(slide, 1)  # Ensure placeholder exists
    
    title.text = f"Content: {topic}"
    
    if body:  # Check if the content placeholder exists
        tf = body.text_frame
        tf.clear()  # Clear any existing content
        p = tf.add_paragraph()
        p.text = content
    else:
        logger.warning(f"No content placeholder found in the second slide of the template {template_file}")

    # Mark 50% completion after generating slides
    progress_bar.progress(50)

    return prs

# Function to generate and add charts to the presentation
def generate_charts_to_presentation(prs, progress_bar):
    """Generates charts and adds them to the presentation."""
    # Generate Pie Chart if selected by the user
    if generate_charts:
        fig, ax = plt.subplots(figsize=(5, 3))
        labels = ['Part A', 'Part B', 'Part C', 'Part D']
        sizes = [15, 30, 45, 10]
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90)
        ax.axis('equal')
        
        chart_path = "chart.png"
        plt.savefig(chart_path)
        plt.close()

        # Add the pie chart to the PowerPoint slide
        slide_layout = prs.slide_layouts[5]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Generated Pie Chart"
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8.5), height=Inches(4.5))

        # Mark 75% completion after generating the pie chart
        progress_bar.progress(75)

    # Generate Bar Chart if selected by the user
    if bar_chart:
        fig = px.bar(
            x=['Category A', 'Category B', 'Category C', 'Category D'],
            y=[10, 20, 30, 40],
            title="Sample Bar Chart"
        )
        chart_path = "bar_chart.html"
        fig.write_html(chart_path)

        # Add the bar chart to the PowerPoint presentation
        slide_layout = prs.slide_layouts[5]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Generated Bar Chart"
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8.5), height=Inches(4.5))

        # Mark 90% completion after generating the bar chart
        progress_bar.progress(90)

    return prs

# Generate content when a topic is provided
if topic and selected_template_name:
    # Show progress bar for content generation
    progress_bar = st.progress(0)
    
    st.subheader(f"Generated Content for '{topic}'")
    with st.spinner("Generating content for the presentation..."):
        content = generate_presentation_content(topic, progress_bar)
        st.markdown(content)

    # Load the selected template file
    template_file_path = os.path.join(templates_folder, f"{selected_template_name}.pptx")
    
    # Create a PowerPoint presentation from the selected template
    prs = create_presentation_from_template(content, topic, template_file_path, progress_bar)

    # Generate and add charts if selected by the user
    prs = generate_charts_to_presentation(prs, progress_bar)

    # Save the presentation and create a download link
    output_file = BytesIO()
    prs.save(output_file)
    output_file.seek(0)

    # Mark 100% completion after everything is done
    progress_bar.progress(100)

    # Provide a download button for the generated PowerPoint file
    st.download_button(
        "Download PowerPoint Presentation",
        data=output_file,
        file_name="generated_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

else:
    st.info("Please enter a topic and select a template to generate the presentation.")
